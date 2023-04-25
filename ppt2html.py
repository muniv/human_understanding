import re
import pptx
from bs4 import BeautifulSoup
from pptx.enum.shapes import MSO_SHAPE_TYPE



DOUBLESPACE_PATTERN  = re.compile(r"\s+")

def ppt2html(ppt_path: str, debug: bool=False):
    prs = pptx.Presentation(ppt_path)
    output = ""

    for _, slide in enumerate(prs.slides):
        # For debugging
        if debug:
            print(f"{_}-slide")
            print("="*20)
        single_slide_output = ""

        # Reordering for sequential access from top to bottom
        sequential_shapes = []

        for idx,shape in enumerate(slide.shapes):
            sequential_shapes.append((shape,shape.top))

        sequential_shapes = sorted(sequential_shapes,key=lambda x : x[1])


        for shape,_ in sequential_shapes:

            # Case 1 : Groupshape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for shape in shape.shapes:
                    if shape.has_text_frame:
                        single_slide_output +=DOUBLESPACE_PATTERN.sub(" ",shape.text_frame.text)+"\n"
                        output +=DOUBLESPACE_PATTERN.sub(" ",shape.text_frame.text)+"\n"

            # Case 2 : Table
            elif shape.has_table:
                soup = BeautifulSoup()
                html_table = soup.new_tag('table')
                for row in shape.table.rows:
                    html_row = soup.new_tag("tr")
                    for cell in row.cells:
                        html_cell = soup.new_tag("td")
                        html_cell.string = cell.text
                        html_row.append(html_cell)
                    html_table.append(html_row)
                single_slide_output+=str(html_table).strip() + "\n"
                output += str(html_table).strip() + "\n"

            # Case 3 : Text box
            elif shape.has_text_frame :
                single_slide_output +=DOUBLESPACE_PATTERN.sub(" ",shape.text_frame.text)+"\n"
                output +=DOUBLESPACE_PATTERN.sub(" ",shape.text_frame.text)+"\n"

        # For debugging
        if debug:
            print(single_slide_output)
            print("="*20)
        output += "\n"

    return output

print(ppt2html('data/test.pptx'))